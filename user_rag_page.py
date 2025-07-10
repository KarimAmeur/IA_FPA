import streamlit as st
import os
import tempfile
import time
import logging
import re
from typing import List
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_community.vectorstores import Chroma
import pandas as pd
import PyPDF2
import docx

# Import pour l'API Mistral (utilise la m√™me classe que app.py)
from mistralai.client import MistralClient

# Import pour PowerPoint
try:
    from pptx import Presentation
except ImportError:
    # Tenter d'installer pptx si non disponible
    import subprocess
    import sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

# Configuration du logging
logging.basicConfig(level=logging.INFO,
                   format='%(asctime)s - %(levelname)s - %(message)s')

# Configuration des APIs - Utilise les secrets Streamlit
try:
    MISTRAL_API_KEY = st.secrets["MISTRAL_API_KEY"]
except:
    MISTRAL_API_KEY = os.getenv("MISTRAL_API_KEY", "")

# CORRECTION MOBILE : Fonction de nettoyage s√©curis√©e
def clean_text_for_mobile(text):
    """Nettoie le texte pour √©viter les erreurs regex sur mobile"""
    if not text or not isinstance(text, str):
        return ""
    
    try:
        # Remplacer les caract√®res probl√©matiques pour mobile
        text = text.replace('\u00a0', ' ')  # Espace ins√©cable
        text = text.replace('\u2019', "'")  # Apostrophe courbe
        text = text.replace('\u201c', '"')  # Guillemet ouvrant
        text = text.replace('\u201d', '"')  # Guillemet fermant
        text = text.replace('\u2013', '-')  # Tiret demi-cadratin
        text = text.replace('\u2014', '-')  # Tiret cadratin
        
        # CORRECTION MOBILE : Nettoyer les URLs SANS regex complexe
        if 'http' in text:
            words = text.split()
            cleaned_words = []
            for word in words:
                if word.startswith(('http://', 'https://', 'www.')):
                    cleaned_words.append('[LIEN]')
                else:
                    cleaned_words.append(word)
            text = ' '.join(cleaned_words)
        
        # CORRECTION MOBILE : Nettoyer les caract√®res de contr√¥le SANS regex complexe
        control_chars = '\x00\x01\x02\x03\x04\x05\x06\x07\x08\x0b\x0c\x0e\x0f\x10\x11\x12\x13\x14\x15\x16\x17\x18\x19\x1a\x1b\x1c\x1d\x1e\x1f\x7f'
        for char in control_chars:
            text = text.replace(char, '')
        
        return text
    except Exception:
        # En cas d'erreur, retourner le texte original sans modification
        return str(text)

class MistralEmbeddings:
    """
    Wrapper LangChain pour Mistral Embed (1024 dims).
    Compatible avec la base vectorielle cr√©√©e par rag_formation.py
    """
    def __init__(self, api_key: str, model: str = "mistral-embed"):
        self.client = MistralClient(api_key=api_key)
        self.model = model

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        embeddings = []
        batch_size = 50
        for i in range(0, len(texts), batch_size):
            batch = texts[i:i+batch_size]
            try:
                # CORRECTION MOBILE : Nettoyer les textes avant embedding
                clean_batch = [clean_text_for_mobile(text) for text in batch]
                resp = self.client.embeddings(model=self.model, input=clean_batch)
                embeddings.extend([d.embedding for d in resp.data])
            except Exception as e:
                st.error(f"Erreur embedding lot {i//batch_size+1}: {e}")
                embeddings.extend([[0.0]*1024 for _ in batch])
        return embeddings

    def embed_query(self, text: str) -> List[float]:
        try:
            # CORRECTION MOBILE : Nettoyer la requ√™te avant embedding
            clean_text = clean_text_for_mobile(text)
            resp = self.client.embeddings(model=self.model, input=[clean_text])
            return resp.data[0].embedding
        except Exception as e:
            st.error(f"Erreur embedding requ√™te: {e}")
            return [0.0]*1024

@st.cache_resource
def get_embedding_model():
    """Charge le mod√®le d'embedding Mistral compatible avec la base vectorielle"""
    try:
        if not MISTRAL_API_KEY:
            st.error("‚ùå Cl√© API Mistral manquante")
            return None
            
        return MistralEmbeddings(api_key=MISTRAL_API_KEY, model="mistral-embed")
        
    except Exception as e:
        st.error(f"‚ùå Erreur lors du chargement du mod√®le d'embedding: {e}")
        return None

def extract_text_from_pdf(file_path: str) -> str:
    """Extraire le texte d'un fichier PDF - VERSION MOBILE SAFE"""
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            full_text = ""
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    # CORRECTION MOBILE : Nettoyer le texte extrait
                    clean_text = clean_text_for_mobile(text)
                    full_text += clean_text + "\n\n"
            return full_text
    except Exception as e:
        logging.error(f"Erreur lors de l'extraction du PDF {file_path}: {e}")
        return ""

def extract_text_from_docx(file_path: str) -> str:
    """Extraire le texte d'un fichier Word - VERSION MOBILE SAFE"""
    try:
        doc = docx.Document(file_path)
        texts = []
        for paragraph in doc.paragraphs:
            if paragraph.text:
                # CORRECTION MOBILE : Nettoyer chaque paragraphe
                clean_para = clean_text_for_mobile(paragraph.text)
                texts.append(clean_para)
        return "\n".join(texts)
    except Exception as e:
        logging.error(f"Erreur lors de l'extraction du DOCX {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    """Extraire le texte d'un fichier PowerPoint - VERSION MOBILE SAFE"""
    try:
        prs = Presentation(file_path)
        full_text = ""
        
        # Extraire le texte de chaque diapositive
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extraire le titre de la diapositive
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text
                if title_text.strip():
                    # CORRECTION MOBILE : Nettoyer le titre
                    clean_title = clean_text_for_mobile(title_text)
                    slide_text.append(f"Titre: {clean_title}")
            
            # Extraire le texte de toutes les formes
            for shape in slide.shapes:
                text_content = ""
                if hasattr(shape, "text") and shape.text.strip():
                    text_content = shape.text
                elif hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    text_content = shape.text_frame.text
                
                if text_content:
                    # CORRECTION MOBILE : Nettoyer le texte de la forme
                    clean_content = clean_text_for_mobile(text_content)
                    slide_text.append(clean_content)
            
            # Ajouter le texte de la diapositive au texte global
            if slide_text:
                full_text += f"Diapositive {i}:\n"
                full_text += "\n".join(slide_text) + "\n\n"
        
        return full_text
    except Exception as e:
        logging.error(f"Erreur lors de l'extraction du PowerPoint {file_path}: {e}")
        return ""

def extract_text_from_xlsx(file_path: str) -> str:
    """Extraire le texte d'un fichier Excel - VERSION MOBILE SAFE"""
    try:
        # Lire tous les onglets
        xls = pd.ExcelFile(file_path)
        full_text = ""

        for sheet_name in xls.sheet_names:
            # Lire la feuille
            df = pd.read_excel(file_path, sheet_name=sheet_name)

            # Convertir toutes les colonnes en cha√Æne et concat√©ner
            sheet_texts = []
            for row in df.values:
                row_text = " ".join(str(val) for val in row if pd.notna(val))
                if row_text.strip():
                    # CORRECTION MOBILE : Nettoyer chaque ligne
                    clean_row = clean_text_for_mobile(row_text)
                    sheet_texts.append(clean_row)

            sheet_text = "\n".join(sheet_texts)
            full_text += f"Feuille {sheet_name}:\n{sheet_text}\n\n"

        return full_text
    except Exception as e:
        logging.error(f"Erreur lors de l'extraction de l'XLSX {file_path}: {e}")
        return ""

def user_rag_page():
    """Page permettant √† l'utilisateur d'ajouter ses propres documents au RAG - VERSION MOBILE SAFE"""
    
    # V√©rification de la cl√© API Mistral
    if not MISTRAL_API_KEY:
        st.error("‚ùå Cl√© API Mistral manquante. Veuillez configurer MISTRAL_API_KEY dans les secrets Streamlit.")
        st.info("üí° Cette fonctionnalit√© n√©cessite une cl√© API Mistral pour vectoriser vos documents avec le m√™me mod√®le d'embedding que la base principale.")
        return
    
    # Banni√®re avec logo et titre - CHARTE EDSET
    st.markdown("""
    <div class="banner">
        <h1>üìÑ edset.</h1>
        <p>importation de documents personnels</p>
    </div>
    """, unsafe_allow_html=True)
    
    # V√©rification de la disponibilit√© des ressources n√©cessaires
    if 'RAG_user' not in st.session_state:
        st.session_state.RAG_user = None
    
    # Colonnes pour l'interface
    left_col, right_col = st.columns([3, 2])
    
    with left_col:
        st.markdown("""
        <div class="scenario-card">
            <h3>üîç Importation de Documents</h3>
            <p>T√©l√©chargez vos documents pour enrichir la base de connaissances personnalis√©e.</p>
            <div class="info-box">
                ‚ÑπÔ∏è Vos documents seront vectoris√©s avec le m√™me mod√®le d'embedding Mistral que la base principale pour assurer la compatibilit√©.
            </div>
        </div>
        """, unsafe_allow_html=True)
        
        # Widget de t√©l√©chargement de fichiers
        uploaded_files = st.file_uploader(
            "S√©lectionnez un ou plusieurs fichiers (PDF, DOCX, PPTX, XLSX)",
            type=["pdf", "docx", "pptx", "ppt", "xlsx", "xls"],
            accept_multiple_files=True,
            help="Les fichiers seront trait√©s puis ajout√©s √† votre RAG personnel"
        )
        
        # Param√®tres pour les chunks - align√©s avec rag_formation.py
        chunk_size = 1024  # M√™me taille que dans rag_formation.py
        chunk_overlap = 100  # M√™me chevauchement que dans rag_formation.py
        
        st.info(f"Les documents seront automatiquement d√©coup√©s en chunks de {chunk_size} caract√®res, avec un chevauchement de {chunk_overlap} caract√®res (param√®tres align√©s avec la base principale).")
        
        # Bouton de traitement
        process_button = st.button(
            "üìä Traiter et vectoriser les documents",
            type="primary",
            use_container_width=True,
            disabled=not uploaded_files
        )
        
        # Traitement des fichiers
        if process_button and uploaded_files:
            with st.status("Traitement des documents en cours...", expanded=True) as status:
                # Cr√©ation d'un dossier temporaire pour enregistrer les fichiers
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Sauvegarde des fichiers dans le dossier temporaire
                    file_paths = []
                    for uploaded_file in uploaded_files:
                        file_path = os.path.join(temp_dir, uploaded_file.name)
                        with open(file_path, "wb") as f:
                            f.write(uploaded_file.getbuffer())
                        file_paths.append(file_path)
                        st.write(f"Fichier sauvegard√© : {uploaded_file.name}")
                    
                    # Extraction du texte des fichiers
                    documents = []
                    
                    for file_path in file_paths:
                        st.write(f"Traitement de {os.path.basename(file_path)}...")
                        
                        # D√©terminer le type de fichier et extraire le texte
                        file_extension = os.path.splitext(file_path)[1].lower()
                        
                        if file_extension == '.pdf':
                            text = extract_text_from_pdf(file_path)
                        elif file_extension == '.docx':
                            text = extract_text_from_docx(file_path)
                        elif file_extension in ['.pptx', '.ppt']:
                            text = extract_text_from_pptx(file_path)
                        elif file_extension in ['.xlsx', '.xls']:
                            text = extract_text_from_xlsx(file_path)
                        else:
                            st.warning(f"Type de fichier non support√© : {file_extension}")
                            continue
                        
                        if not text.strip():
                            st.warning(f"Aucun texte extrait de {os.path.basename(file_path)}")
                            continue
                        
                        # CORRECTION MOBILE : Nettoyer le texte extrait
                        text = clean_text_for_mobile(text)
                        
                        # Limitation de taille comme dans rag_formation.py
                        if len(text) > 500_000:
                            st.warning(f"Tronqu√© {os.path.basename(file_path)} ({len(text)}‚Üí500000)")
                            text = text[:500_000]
                        
                        st.write(f"‚úÖ Texte extrait ({len(text)} caract√®res)")
                        
                        # D√©coupage du texte en chunks avec les m√™mes param√®tres que rag_formation.py
                        text_splitter = RecursiveCharacterTextSplitter(
                            chunk_size=chunk_size,
                            chunk_overlap=chunk_overlap,
                            length_function=len,
                            separators=["\n\n", "\n", ".", "!", "?", ";", ":", " ", ""]
                        )
                        
                        chunks = text_splitter.split_text(text)
                        st.write(f"üî™ Texte d√©coup√© en {len(chunks)} chunks")
                        
                        # Cr√©ation des documents Langchain avec filtrage comme dans rag_formation.py
                        filename = os.path.basename(file_path)
                        chunks_added = 0
                        for i, chunk in enumerate(chunks):
                            if len(chunk.strip()) < 100:  # M√™me filtrage que rag_formation.py
                                continue
                            
                            # CORRECTION MOBILE : Nettoyer chaque chunk
                            clean_chunk = clean_text_for_mobile(chunk)
                            
                            documents.append(Document(
                                page_content=clean_chunk,
                                metadata={
                                    "source": file_path,
                                    "filename": filename,
                                    "chunk": i,
                                    "type": file_extension[1:],  # Type sans le point initial
                                    "size": len(clean_chunk)
                                }
                            ))
                            chunks_added += 1
                        
                        st.write(f"   ‚Üí {chunks_added} chunks ajout√©s (apr√®s filtrage)")
                    
                    # Cr√©ation ou mise √† jour de la base vectorielle
                    if documents:
                        st.write(f"üß† Vectorisation de {len(documents)} documents avec Mistral Embed...")
                        
                        try:
                            # R√©cup√©rer le mod√®le d'embedding Mistral
                            embeddings = get_embedding_model()
                            
                            if embeddings is None:
                                status.update(label="Erreur: Mod√®le d'embedding non disponible", state="error")
                                return
                            
                            # D√©terminer s'il faut cr√©er une nouvelle base ou mettre √† jour l'existante
                            if st.session_state.RAG_user is None:
                                # Cr√©ation d'une nouvelle base vectorielle
                                persist_directory = "chroma_db_user"
                                st.session_state.RAG_user = Chroma.from_documents(
                                    documents=documents,
                                    embedding=embeddings,
                                    persist_directory=persist_directory
                                )
                                st.write(f"üéâ Nouvelle base vectorielle cr√©√©e avec {len(documents)} documents")
                            else:
                                # Mise √† jour de la base existante
                                st.session_state.RAG_user.add_documents(documents=documents)
                                st.write(f"üîÑ Base vectorielle mise √† jour avec {len(documents)} nouveaux documents")
                            
                            status.update(label="Traitement termin√© avec succ√®s!", state="complete", expanded=False)
                            
                            # Indiquer le nombre total de documents dans la base
                            try:
                                collection = st.session_state.RAG_user._collection
                                total_docs = collection.count()
                                st.success(f"Votre RAG personnel contient maintenant {total_docs} chunks de documents")
                            except:
                                st.success("Vectorisation termin√©e avec succ√®s!")
                            
                            st.rerun()
                        except Exception as e:
                            status.update(label=f"Erreur lors de la vectorisation: {str(e)}", state="error")
                            st.error(f"Erreur de vectorisation: {str(e)}")
                    else:
                        status.update(label="Aucun document valide √† vectoriser", state="error")
                        st.warning("Aucun document valide n'a pu √™tre trait√©")
    
    
    with right_col:
        st.markdown("""
        <div class="scenario-card">
            <h3>‚ÑπÔ∏è Informations sur votre RAG personnel</h3>
        </div>
        """, unsafe_allow_html=True)
        
        # Affichage des informations sur le RAG personnel
        if st.session_state.RAG_user is not None:
            try:
                # R√©cup√©rer les informations sur la base vectorielle
                docs_dict = st.session_state.RAG_user.get()
                total_docs = len(docs_dict['documents'])
            
                # R√©cup√©rer les m√©tadonn√©es des documents
                sources = {}
                file_types = {}
                
                for metadata in docs_dict['metadatas']:
                    filename = metadata.get('filename', 'Inconnu')
                    filetype = metadata.get('type', 'Inconnu')
                    
                    # Compter par nom de fichier
                    if filename in sources:
                        sources[filename] += 1
                    else:
                        sources[filename] = 1
                    
                    # Compter par type de fichier
                    if filetype in file_types:
                        file_types[filetype] += 1
                    else:
                        file_types[filetype] = 1
                
                st.markdown(f"""
                <div class="info-box">
                    <p><strong>üìä Statistiques de votre RAG personnel:</strong></p>
                    <ul>
                        <li><strong>Nombre total de chunks:</strong> {total_docs}</li>
                        <li><strong>Nombre de fichiers sources:</strong> {len(sources)}</li>
                        <li><strong>Mod√®le d'embedding:</strong> Mistral Embed (1024 dims)</li>
                    </ul>
                </div>
                """, unsafe_allow_html=True)
                
                # Afficher la r√©partition par type de fichier
                st.markdown("<p><strong>üìä R√©partition par type de fichier:</strong></p>", unsafe_allow_html=True)
                for filetype, count in file_types.items():
                    file_emoji = {
                        'pdf': 'üìÑ', 
                        'docx': 'üìù', 
                        'pptx': 'üìä',
                        'ppt': 'üìä',
                        'xlsx': 'üìà',
                        'xls': 'üìà'
                    }.get(filetype, 'üìÅ')
                    
                    st.markdown(f"- {file_emoji} **{filetype.upper()}**: {count} chunks")
                
                # Afficher la liste des fichiers
                st.markdown("<p><strong>üìÅ Fichiers vectoris√©s:</strong></p>", unsafe_allow_html=True)
                for filename, count in sources.items():
                    st.markdown(f"- **{filename}**: {count} chunks")
                
                # Option pour tester le RAG
                st.markdown("""
                <div class="scenario-card" style="margin-top: 20px;">
                    <h3>üîç Tester votre RAG personnel</h3>
                </div>
                """, unsafe_allow_html=True)
                
                test_query = st.text_input(
                    "Saisissez une requ√™te de test",
                    placeholder="Exemple: Quels sont les points cl√©s abord√©s dans mes documents?",
                    help="Cette requ√™te sera utilis√©e pour rechercher dans vos documents"
                )
                
                if st.button("üîé Rechercher", use_container_width=True) and test_query:
                    # Effectuer la recherche
                    with st.spinner("Recherche en cours..."):
                        try:
                            # CORRECTION MOBILE : Nettoyer la requ√™te de test
                            clean_query = clean_text_for_mobile(test_query)
                            
                            # Effectuer la recherche similaire
                            results = st.session_state.RAG_user.similarity_search_with_score(
                                query=clean_query,
                                k=3  # Nombre de r√©sultats √† afficher
                            )
                            
                            if results:
                                st.markdown("<p><strong>üìë R√©sultats de la recherche:</strong></p>", unsafe_allow_html=True)
                                for i, (doc, score) in enumerate(results, 1):
                                    file_type = doc.metadata.get('type', 'inconnu')
                                    file_emoji = {
                                        'pdf': 'üìÑ', 
                                        'docx': 'üìù', 
                                        'pptx': 'üìä',
                                        'ppt': 'üìä',
                                        'xlsx': 'üìà',
                                        'xls': 'üìà'
                                    }.get(file_type, 'üìÅ')
                                    
                                    # CORRECTION MOBILE : Nettoyer le contenu affich√©
                                    clean_content = clean_text_for_mobile(doc.page_content)
                                    clean_filename = clean_text_for_mobile(doc.metadata.get('filename', 'Inconnu'))
                                    
                                    with st.expander(f"{file_emoji} R√©sultat {i} - Score: {score:.4f} - Source: {clean_filename}"):
                                        st.markdown(f"**Extrait du document:**\n\n{clean_content}")
                            else:
                                st.info("Aucun r√©sultat correspondant √† votre requ√™te")
                        except Exception as e:
                            st.error(f"Erreur lors de la recherche: {str(e)}")
            except Exception as e:
                st.error(f"Erreur lors de la r√©cup√©ration des informations sur le RAG: {str(e)}")
        else:
            st.info("Aucun document n'a encore √©t√© import√© dans votre RAG personnel")
            
            st.markdown("""
            <div class="info-box" style="margin-top: 20px;">
                <p><strong>‚ÑπÔ∏è Comment utiliser cette fonctionnalit√©:</strong></p>
                <ol>
                    <li>T√©l√©chargez vos documents (PDF, Word, PowerPoint, Excel)</li>
                    <li>Cliquez sur "Traiter et vectoriser les documents"</li>
                    <li>Une fois vos documents vectoris√©s, vous pourrez tester la recherche</li>
                </ol>
                <p><strong>üîß Compatibilit√©:</strong> Vos documents seront vectoris√©s avec le m√™me mod√®le Mistral Embed (1024 dimensions) que la base principale, garantissant une compatibilit√© parfaite.</p>
                <p><strong>üíæ Stockage:</strong> Vos documents sont stock√©s uniquement dans votre session Streamlit actuelle et ne seront pas conserv√©s apr√®s la fermeture de l'application.</p>
            </div>
            """, unsafe_allow_html=True)
        
        # Ajouter un bouton pour effacer le RAG personnel
        if st.session_state.RAG_user is not None:
            st.markdown("""
            <div class="scenario-card" style="margin-top: 20px;">
                <h3>üóëÔ∏è Gestion du RAG personnel</h3>
            </div>
            """, unsafe_allow_html=True)
            
            if st.button("üóëÔ∏è Effacer mon RAG personnel", type="secondary", use_container_width=True):
                # Demander confirmation
                st.warning("‚ö†Ô∏è Cette action supprimera d√©finitivement tous vos documents vectoris√©s.")
                confirm = st.checkbox("Confirmer la suppression de tous mes documents")
                
                if confirm and st.button("‚úÖ Confirmer la suppression", type="secondary"):
                    try:
                        # Supprimer le dossier de persistance s'il existe
                        import shutil
                        if os.path.exists("chroma_db_user"):
                            shutil.rmtree("chroma_db_user")
                        
                        # R√©initialiser la session
                        st.session_state.RAG_user = None
                        st.success("Votre RAG personnel a √©t√© effac√© avec succ√®s!")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Erreur lors de la suppression: {str(e)}")